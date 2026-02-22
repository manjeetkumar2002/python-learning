import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import networkx as nx


# 1. Graph Image Generate karne ka function
def generate_graph_image():
    G = nx.Graph()
    edges = [('0', '1'), ('0', '2'), ('1', '3'), ('1', '4'), ('2', '5'), ('5', '6'), ('5', '7')]
    G.add_edges_from(edges)

    pos = {'0': (0, 2), '1': (-1, 1), '2': (1, 1), '3': (-1.5, 0), '4': (-0.5, 0), '5': (1, 0), '6': (0.5, -1),
           '7': (1.5, -1)}

    plt.figure(figsize=(6, 4))
    nx.draw(G, pos, with_labels=True, node_color='skyblue', node_size=2000, edge_color='black', linewidths=2,
            font_size=15)
    plt.title("Experimental Graph Structure")
    plt.savefig("graph_dataset.png")
    plt.close()


# 2. PPT Banane ka function
def create_advanced_ppt():
    generate_graph_image()
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Performance Analysis of BFS and DFS"
    slide.placeholders[
        1].text = "Applied AI - Project Presentation\n\nName: [Tera Naam]\nRoll No: [Tera Roll No]\nGuide: [Mam's Name]"

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agenda (Vishay Suchi)"
    points = ["Introduction to Traversal", "BFS vs DFS Overview", "Dataset/Graph Description",
              "Code Implementation", "Performance Metrics", "Comparative Analysis", "Conclusion"]
    for p in points:
        slide.placeholders[1].text_frame.add_paragraph().text = p

    # --- Slide 3: Intro ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction to Graph Traversal"
    tf = slide.placeholders[1].text_frame
    tf.text = "Graph traversal ka matlab hai har node ko exactly ek baar visit karna."
    tf.add_paragraph().text = "AI Applications: Pathfinding, Web Crawling, Social Network Analysis."

    # --- Slide 4 & 5: BFS & DFS Details ---
    for title, strategy, ds in [("What is BFS?", "Level-by-Level (Breadth-wise)", "Queue (FIFO)"),
                                ("What is DFS?", "Go Deep (Branch-wise)", "Stack / Recursion")]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.text = f"Strategy: {strategy}"
        tf.add_paragraph().text = f"Data Structure: {ds}"

    # --- Slide 6: Dataset (With Generated Image) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank title layout
    slide.shapes.title.text = "Our Dataset (The Experimental Graph)"
    slide.shapes.add_picture("graph_dataset.png", Inches(1), Inches(2), width=Inches(8))

    # --- Slide 9 & 10: Code Snippets (Font chota rakha hai taaki fit ho jaye) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "BFS/DFS Implementation (Python)"
    code_box = slide.placeholders[1].text_frame
    code_box.text = "def bfs(graph, start):\n    visited = set()\n    queue = deque([start])\n    ..."
    p = code_box.add_paragraph()
    p.text = "\ndef dfs(graph, start, visited=None):\n    if visited is None: visited = set()\n    ..."
    p.font.size = Pt(14)

    # --- Slide 13: Time Complexity Table ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Performance Metrics: Time Complexity"
    table = slide.shapes.add_table(3, 3, Inches(0.5), Inches(2), Inches(9), Inches(2)).table
    headers = ["Algorithm", "Time Complexity", "Reason"]
    data = [("BFS", "O(V + E)", "Every node & edge visited once"),
            ("DFS", "O(V + E)", "Every node & edge visited once")]
    for i, h in enumerate(headers): table.cell(0, i).text = h
    for r, (a, t, res) in enumerate(data, 1):
        table.cell(r, 0).text = a
        table.cell(r, 1).text = t
        table.cell(r, 2).text = res

    # --- Slide 16: Half-Graph Analysis (Special Slide) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Half-Graph Performance Analysis"
    tf = slide.placeholders[1].text_frame
    tf.text = "Observation based on splitting the graph:"
    tf.add_paragraph().text = "• BFS: Junction nodes (Node 2) se start karke dono parts ko balance tarike se explore kiya."
    tf.add_paragraph().text = "• DFS: Pehle Part A ki depth (1->3->4) khatam ki, phir Part B par gaya."

    # --- Slide 18: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You!"
    slide.placeholders[1].text = "Questions? \nReferences: GeeksforGeeks, Applied AI Textbooks"

    prs.save("BFS_DFS_Performance_Final.pptx")
    print("Bhai, PPT 'BFS_DFS_Performance_Final.pptx' ke naam se save ho gayi hai!")


create_advanced_ppt()